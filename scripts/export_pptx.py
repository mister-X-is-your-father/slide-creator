#!/usr/bin/env python3
"""
export_pptx.py - MarpスライドをPowerPointにエクスポート

使い方:
    python export_pptx.py slide.md --output 3_export/pptx/
    python export_pptx.py slide.md --template assets/templates/corporate.potx
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# プロジェクトルート
PROJECT_ROOT = Path(__file__).parent.parent

# デフォルト出力ディレクトリ
DEFAULT_OUTPUT = PROJECT_ROOT / "3_export" / "pptx"


def parse_marp_slides(content: str) -> list:
    """Marpスライドをパースしてセクションのリストにする"""
    # フロントマターを除去
    content = re.sub(r'^---\n.*?\n---\n', '', content, flags=re.DOTALL)

    # スライド区切りで分割
    slides = re.split(r'\n---\n', content)

    parsed = []
    for slide in slides:
        slide = slide.strip()
        if not slide:
            continue

        section = {
            'type': 'content',
            'title': '',
            'content': [],
            'raw': slide
        }

        # タイトルスライド検出
        if '<div class="title-slide">' in slide:
            section['type'] = 'title'
            match = re.search(r'#\s+(.+)', slide)
            if match:
                section['title'] = match.group(1).strip()
            clean = re.sub(r'<[^>]+>', '', slide)
            clean = re.sub(r'#\s+.+', '', clean)
            section['content'] = [line.strip() for line in clean.split('\n') if line.strip()]

        # セクションヘッダー検出
        elif '<div class="section-header">' in slide:
            section['type'] = 'section'
            match = re.search(r'#\s+(.+)', slide)
            if match:
                section['title'] = match.group(1).strip()

        # 通常スライド
        else:
            match = re.search(r'##\s+(.+)', slide)
            if match:
                section['title'] = match.group(1).strip()

            items = re.findall(r'^[-*]\s+(.+)$', slide, re.MULTILINE)
            if items:
                section['content'] = items
            else:
                clean = re.sub(r'<[^>]+>', '', slide)
                clean = re.sub(r'##?\s+.+', '', clean)
                section['content'] = [line.strip() for line in clean.split('\n') if line.strip()]

        parsed.append(section)

    return parsed


def hex_to_rgb(hex_color: str) -> RgbColor:
    """16進数カラーコードをRgbColorに変換"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RgbColor(r, g, b)


def create_pptx(slides: list, output_file: Path, template: Path = None):
    """PowerPointファイルを作成"""

    # テンプレートがあれば使用
    if template and template.exists():
        prs = Presentation(str(template))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

    # スライドレイアウト
    # 0: Title Slide, 1: Title and Content, 5: Blank, 6: Title Only
    try:
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]
        section_layout = prs.slide_layouts[5]  # Blank
        blank_layout = prs.slide_layouts[6]    # Title Only
    except IndexError:
        # カスタムテンプレートの場合
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[0]
        section_layout = prs.slide_layouts[0]
        blank_layout = prs.slide_layouts[0]

    for section in slides:
        if section['type'] == 'title':
            # タイトルスライド
            slide = prs.slides.add_slide(title_layout)
            if slide.shapes.title:
                slide.shapes.title.text = section['title']

            # サブタイトル
            if section['content'] and len(slide.placeholders) > 1:
                try:
                    subtitle = slide.placeholders[1]
                    subtitle.text = '\n'.join(section['content'])
                except (KeyError, IndexError):
                    pass

        elif section['type'] == 'section':
            # セクションヘッダー
            slide = prs.slides.add_slide(section_layout)
            # テキストボックスを追加
            left = Inches(1)
            top = Inches(3)
            width = Inches(11)
            height = Inches(1.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.paragraphs[0].text = section['title']
            tf.paragraphs[0].font.size = Pt(44)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        else:
            # 通常コンテンツスライド
            slide = prs.slides.add_slide(content_layout)

            # タイトル
            if slide.shapes.title and section['title']:
                slide.shapes.title.text = section['title']

            # コンテンツ
            if section['content']:
                # コンテンツプレースホルダーを探す
                body_shape = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # Body placeholder
                        body_shape = shape
                        break

                if body_shape:
                    tf = body_shape.text_frame
                    tf.clear()  # 既存テキストをクリア

                    for i, item in enumerate(section['content']):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
                        p.font.size = Pt(18)
                else:
                    # プレースホルダーがない場合はテキストボックス
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(11)
                    height = Inches(5)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame

                    for i, item in enumerate(section['content']):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(18)

    # 保存
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    return output_file


def list_slides(slides_dir: Path) -> list:
    """スライドディレクトリ内のファイル一覧"""
    if not slides_dir.exists():
        return []
    return sorted(slides_dir.glob("*.md"))


def interactive_select(slides: list) -> Path:
    """対話的にスライドを選択"""
    print("\n利用可能なスライド:\n")
    for i, slide in enumerate(slides, 1):
        print(f"  {i}. {slide.name}")

    print()
    while True:
        try:
            choice = input("番号を選択 (q で終了): ").strip()
            if choice.lower() == "q":
                sys.exit(0)
            idx = int(choice) - 1
            if 0 <= idx < len(slides):
                return slides[idx]
            print("無効な番号です")
        except ValueError:
            print("数字を入力してください")


def main():
    parser = argparse.ArgumentParser(
        description="MarpスライドをPowerPointにエクスポート"
    )
    parser.add_argument(
        "input",
        nargs="?",
        help="入力ファイルパス (省略時は対話選択)"
    )
    parser.add_argument(
        "--output", "-o",
        help="出力先ディレクトリ (デフォルト: 3_export/pptx/)"
    )
    parser.add_argument(
        "--template", "-t",
        help="PowerPointテンプレート (.potx/.pptx)"
    )

    args = parser.parse_args()

    # python-pptx 確認
    if not PPTX_AVAILABLE:
        print("エラー: python-pptx がインストールされていません")
        print("インストール: pip install python-pptx")
        sys.exit(1)

    # 入力ファイル
    if args.input:
        input_file = Path(args.input)
        if not input_file.exists():
            print(f"エラー: ファイルが見つかりません: {input_file}")
            sys.exit(1)
    else:
        slides_dir = PROJECT_ROOT / "2_create"
        slides = list_slides(slides_dir)
        if not slides:
            print("エラー: 2_create/ にスライドがありません")
            sys.exit(1)
        input_file = interactive_select(slides)

    # 出力先
    if args.output:
        output_dir = Path(args.output)
    else:
        output_dir = DEFAULT_OUTPUT

    output_file = output_dir / f"{input_file.stem}.pptx"

    # テンプレート
    template = Path(args.template) if args.template else None
    if template and not template.exists():
        print(f"警告: テンプレートが見つかりません: {template}")
        template = None

    # スライド読み込み・パース
    with open(input_file, "r", encoding="utf-8") as f:
        content = f.read()

    slides = parse_marp_slides(content)
    print(f"\n{len(slides)} スライドを検出")

    # エクスポート実行
    print(f"\nPowerPoint作成中...")
    try:
        result = create_pptx(
            slides=slides,
            output_file=output_file,
            template=template
        )
        print(f"\n✓ 作成完了: {result}")
    except Exception as e:
        print(f"\nエラー: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
